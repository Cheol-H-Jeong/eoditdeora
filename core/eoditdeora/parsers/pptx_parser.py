"""PPTX parser via python-pptx."""

from __future__ import annotations

from pathlib import Path

from eoditdeora.parsers.base import Block, ParsedDoc, ParseResult, ParserError
from eoditdeora.parsers.registry import register
from eoditdeora.utils.paths_util import display_path


class PptxParser:
    name = "pptx_python_pptx"
    supported_extensions = ("pptx", "pptm")
    fidelity = 4

    def can_parse(self, path: Path) -> bool:
        return path.suffix.lower().lstrip(".") in self.supported_extensions

    def parse(self, path: Path, *, doc_id: str) -> ParseResult:
        try:
            from pptx import Presentation  # type: ignore[import-not-found]
        except ImportError as e:
            raise ParserError("python-pptx not available") from e

        try:
            prs = Presentation(str(path))
        except Exception as e:  # noqa: BLE001
            raise ParserError(f"pptx_open_failed: {e}") from e

        blocks: list[Block] = []
        for idx, slide in enumerate(prs.slides, start=1):
            texts: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        runs = "".join(run.text for run in para.runs)
                        if runs.strip():
                            texts.append(runs.strip())
                if getattr(shape, "has_table", False):
                    for row in shape.table.rows:
                        cells = "\t".join((c.text or "").strip() for c in row.cells)
                        if cells.strip():
                            texts.append(cells)
            title = _slide_title(slide)
            if title:
                blocks.append(
                    Block(type="heading", text=title, level=2, page=idx, meta={"slide": idx})
                )
            body = "\n".join(t for t in texts if t and t != title)
            if body.strip():
                blocks.append(Block(type="paragraph", text=body, page=idx, meta={"slide": idx}))

        try:
            props = prs.core_properties
            metadata = {
                "author": props.author,
                "created_at": props.created.isoformat() if props.created else None,
                "modified_at": props.modified.isoformat() if props.modified else None,
                "application": "Microsoft PowerPoint",
                "title": props.title,
                "page_count": len(prs.slides),
            }
        except Exception:  # noqa: BLE001
            metadata = {"page_count": len(prs.slides)}

        doc = ParsedDoc(
            doc_id=doc_id,
            source_path=str(path),
            source_path_display=display_path(path),
            format="pptx",
            parser=self.name,
            fidelity=self.fidelity,
            blocks=blocks,
            metadata={k: v for k, v in metadata.items() if v is not None},
        )
        return ParseResult(doc=doc)


def _slide_title(slide) -> str:  # type: ignore[no-untyped-def]
    try:
        title_shape = slide.shapes.title
        if title_shape and title_shape.has_text_frame:
            return (title_shape.text_frame.text or "").strip()
    except AttributeError:
        pass
    return ""


register(PptxParser())
