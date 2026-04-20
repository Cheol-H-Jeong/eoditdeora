"""Test fixture helpers that synthesize real documents at runtime.

Keeps the repo free of binary test artifacts. Every parser gets a
"synthesize + parse + assert" round-trip.
"""

from __future__ import annotations

import zipfile
from pathlib import Path


def make_docx(path: Path, paragraphs: list[tuple[str, str]]) -> Path:
    """Build a minimal .docx with (style_name, text) blocks.

    style_name should be something python-docx recognizes: "Title",
    "Heading 1", "Normal", "List Bullet", "Quote".
    """
    import docx

    d = docx.Document()
    for style, text in paragraphs:
        d.add_paragraph(text, style=style)
    path.parent.mkdir(parents=True, exist_ok=True)
    d.save(str(path))
    return path


def make_xlsx(path: Path, sheets: dict[str, list[list[object]]]) -> Path:
    from openpyxl import Workbook

    wb = Workbook()
    default = wb.active
    wb.remove(default)
    for name, rows in sheets.items():
        ws = wb.create_sheet(name)
        for row in rows:
            ws.append(row)
    path.parent.mkdir(parents=True, exist_ok=True)
    wb.save(str(path))
    return path


def make_pptx(path: Path, slides: list[tuple[str, str]]) -> Path:
    from pptx import Presentation

    prs = Presentation()
    layout = prs.slide_layouts[1]  # Title + Content
    for title, body in slides:
        slide = prs.slides.add_slide(layout)
        if slide.shapes.title:
            slide.shapes.title.text = title
        placeholders = [ph for ph in slide.placeholders if ph.placeholder_format.idx != 0]
        if placeholders:
            placeholders[0].text = body
    path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(path))
    return path


def make_hwpx(path: Path, title: str, paragraphs: list[str]) -> Path:
    """Synthesize a minimal valid HWPX (zip with OOXML-style XML).

    We build just enough of the HWPX 2011 schema that our native parser
    can find the section and enumerate paragraphs. The important parts:
      * a mimetype entry at offset 0
      * a Contents/content.hpf OPF with dc:title and dc:creator
      * a Contents/section0.xml with `hp:p` paragraphs holding `hp:t` runs
    """
    path.parent.mkdir(parents=True, exist_ok=True)

    ns_p = "http://www.hancom.co.kr/hwpml/2011/paragraph"
    ns_s = "http://www.hancom.co.kr/hwpml/2011/section"
    ns_dc = "http://purl.org/dc/elements/1.1/"

    # OPF content
    opf_items = "".join(
        f'<dc:title>{title}</dc:title>'
        f'<dc:creator>tester</dc:creator>'
    )
    content_hpf = (
        f'<?xml version="1.0" encoding="UTF-8"?>'
        f'<opf:package xmlns:opf="http://www.idpf.org/2007/opf/" '
        f'xmlns:dc="{ns_dc}" unique-identifier="bookid">'
        f'<opf:metadata>{opf_items}</opf:metadata>'
        f'</opf:package>'
    )

    # Section with paragraphs
    paras_xml = "".join(
        f'<hp:p styleIDRef="0"><hp:run><hp:t>{p}</hp:t></hp:run></hp:p>'
        for p in paragraphs
    )
    section_xml = (
        f'<?xml version="1.0" encoding="UTF-8"?>'
        f'<hs:sec xmlns:hs="{ns_s}" xmlns:hp="{ns_p}">'
        f'{paras_xml}'
        f'</hs:sec>'
    )

    mimetype = b"application/hwp+zip"

    with zipfile.ZipFile(path, "w", compression=zipfile.ZIP_DEFLATED) as zf:
        # mimetype must be the first entry and stored (uncompressed)
        info = zipfile.ZipInfo("mimetype")
        info.compress_type = zipfile.ZIP_STORED
        zf.writestr(info, mimetype)
        zf.writestr("Contents/content.hpf", content_hpf)
        zf.writestr("Contents/section0.xml", section_xml)
    return path


def make_pdf(path: Path, lines: list[str]) -> Path:
    """Build a minimal PDF with extractable text using pypdfium2-friendly
    bytes hand-rolled. pypdfium2 can parse it.

    This writes a tiny hand-authored PDF containing one page with a
    Helvetica text stream. The cat-n encoding is kept ASCII so the
    string stays stable across platforms.
    """
    # Use ReportLab-free approach: construct PDF manually. For simplicity,
    # rely on reportlab if installed; otherwise build a trivial one-page
    # PDF using the standard library. We prefer a synthesized PDF over a
    # static one so we do not ship binaries in-repo.
    try:
        from pypdf import PdfWriter  # type: ignore[import-not-found]
    except ImportError:  # pragma: no cover
        pass

    # Manual minimal PDF construction — produces a parseable, text-bearing
    # single-page document. Not pretty, but round-trips through pdfplumber.
    def _pdf_escape(s: str) -> str:
        return s.replace("\\", "\\\\").replace("(", "\\(").replace(")", "\\)")

    stream_lines = []
    y = 760
    stream_lines.append("BT")
    stream_lines.append("/F1 14 Tf")
    for line in lines:
        stream_lines.append(f"1 0 0 1 72 {y} Tm")
        stream_lines.append(f"({_pdf_escape(line)}) Tj")
        y -= 18
    stream_lines.append("ET")
    stream_content = "\n".join(stream_lines).encode("latin-1")

    objects = []
    def add(obj: bytes) -> int:
        objects.append(obj)
        return len(objects)

    # Object 1: Catalog
    catalog = b"<< /Type /Catalog /Pages 2 0 R >>"
    # Object 2: Pages
    pages = b"<< /Type /Pages /Count 1 /Kids [3 0 R] >>"
    # Object 3: Page
    page = (
        b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] "
        b"/Contents 4 0 R /Resources << /Font << /F1 5 0 R >> >> >>"
    )
    # Object 4: Content stream
    stream_obj = (
        b"<< /Length " + str(len(stream_content)).encode() + b" >>\nstream\n"
        + stream_content + b"\nendstream"
    )
    # Object 5: Font
    font = b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>"

    buf = bytearray(b"%PDF-1.4\n%\xe2\xe3\xcf\xd3\n")
    offsets = [0]
    for i, body in enumerate([catalog, pages, page, stream_obj, font], start=1):
        offsets.append(len(buf))
        buf += f"{i} 0 obj\n".encode() + body + b"\nendobj\n"
    xref_off = len(buf)
    buf += b"xref\n0 6\n"
    buf += b"0000000000 65535 f \n"
    for off in offsets[1:]:
        buf += f"{off:010d} 00000 n \n".encode()
    buf += b"trailer\n<< /Size 6 /Root 1 0 R >>\nstartxref\n"
    buf += f"{xref_off}\n%%EOF".encode()

    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_bytes(bytes(buf))
    return path
